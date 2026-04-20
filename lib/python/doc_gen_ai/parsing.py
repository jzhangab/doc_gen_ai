import io
import logging

logger = logging.getLogger(__name__)


def _get_uploaded_files(upload_widget):
    """Return [{name, content_bytes}] from a FileUpload widget.
    Handles both ipywidgets 7.x (dict) and 8.x (tuple of UploadedFile)."""
    value = upload_widget.value
    result = []
    if isinstance(value, dict):
        for fname, info in value.items():
            content = info.get("content", b"")
            result.append({"name": fname, "content_bytes": bytes(content)})
    else:
        for f in value:
            name = getattr(f, "name", None) or f.get("name", "file")
            content = getattr(f, "content", None) or f.get("content", b"")
            result.append({"name": name, "content_bytes": bytes(content)})
    return result


def extract_text(filename: str, content_bytes: bytes) -> str:
    name = filename.lower()
    try:
        if name.endswith((".docx", ".doc")):
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        if name.endswith(".pdf"):
            import fitz
            doc = fitz.open(stream=content_bytes, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)

        if name.endswith((".xlsx", ".xls")):
            import openpyxl
            wb = openpyxl.load_workbook(
                io.BytesIO(content_bytes), read_only=True, data_only=True
            )
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join("" if c is None else str(c) for c in row)
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts)

        if name.endswith((".pptx", ".ppt")):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)

        return content_bytes.decode("utf-8", errors="replace")

    except Exception as exc:
        logger.error("extract_text failed for %s: %s", filename, exc)
        return f"[Error reading {filename}: {exc}]"
