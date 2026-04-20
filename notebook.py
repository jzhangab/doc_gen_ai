"""
doc_gen_ai — V&V Document Generator (Dataiku Notebook, Python 3.9)

Paste the contents of each cell into a separate Dataiku notebook cell and
run them in order. The full UI appears after the last cell runs.

Required packages (install once in your Dataiku code env):
    python-docx  PyMuPDF  openpyxl  python-pptx  ipywidgets
"""

# ══════════════════════════════════════════════════════════════════════════════
# CELL 1 — Imports & Configuration
# ══════════════════════════════════════════════════════════════════════════════

import base64
import io
import json
import logging
import threading
import traceback
import uuid
from datetime import datetime

import ipywidgets as widgets
from IPython.display import HTML, Javascript, clear_output, display

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

DEFAULT_LLM_CONNECTION_ID = "azureopenai:Azure-OpenAi:gpt-5.2"

DOC_TYPES = [
    "Administration Guide",
    "Infrastructure Configuration Specification",
    "Installation Checklist Protocol",
    "Requirement and Design Specification",
    "System Support Plan",
    "System Test Protocol",
    "User Guide",
    "Verification Plan",
]

_MAX_CHARS = 12000

print("✓ Cell 1 complete — imports loaded")


# ══════════════════════════════════════════════════════════════════════════════
# CELL 2 — Document Parsing Utilities
# ══════════════════════════════════════════════════════════════════════════════

def _get_uploaded_files(upload_widget):
    """Return [{name, content_bytes}] from a FileUpload widget.
    Handles both ipywidgets 7.x (dict) and 8.x (tuple of UploadedFile)."""
    value = upload_widget.value
    result = []
    if isinstance(value, dict):
        for fname, info in value.items():
            content = info.get("content", b"")
            result.append({"name": fname, "content_bytes": bytes(content)})
    else:
        for f in value:
            name = getattr(f, "name", None) or f.get("name", "file")
            content = getattr(f, "content", None) or f.get("content", b"")
            result.append({"name": name, "content_bytes": bytes(content)})
    return result


def extract_text(filename: str, content_bytes: bytes) -> str:
    """Extract plain text from a document given its filename and raw bytes."""
    name = filename.lower()
    try:
        if name.endswith((".docx", ".doc")):
            from docx import Document
            doc = Document(io.BytesIO(content_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        if name.endswith(".pdf"):
            import fitz
            doc = fitz.open(stream=content_bytes, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)

        if name.endswith((".xlsx", ".xls")):
            import openpyxl
            wb = openpyxl.load_workbook(
                io.BytesIO(content_bytes), read_only=True, data_only=True
            )
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    line = "\t".join("" if c is None else str(c) for c in row)
                    if line.strip():
                        parts.append(line)
            return "\n".join(parts)

        if name.endswith((".pptx", ".ppt")):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content_bytes))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)

        return content_bytes.decode("utf-8", errors="replace")

    except Exception as exc:
        logger.error("extract_text failed for %s: %s", filename, exc)
        return f"[Error reading {filename}: {exc}]"


print("✓ Cell 2 complete — document parsing ready")


# ══════════════════════════════════════════════════════════════════════════════
# CELL 3 — LLM Pipeline
# ══════════════════════════════════════════════════════════════════════════════

def _trunc(text: str) -> str:
    return text[:_MAX_CHARS] + "\n...[truncated]" if len(text) > _MAX_CHARS else text


def _llm_call(messages: list, connection_id: str = None) -> str:
    import dataiku
    conn = connection_id or DEFAULT_LLM_CONNECTION_ID
    api_client = dataiku.api_client()
    project = api_client.get_project(dataiku.default_project_key())
    llm = project.get_llm(conn)
    completion = llm.new_completion()
    for msg in messages:
        completion.with_message(msg["content"], msg["role"])
    try:
        completion.with_max_output_tokens(8192)
    except Exception:
        pass
    return completion.execute().text


def _llm_json(messages: list, connection_id: str = None) -> dict:
    raw = _llm_call(messages, connection_id=connection_id).strip()
    if raw.startswith("```"):
        raw = "\n".join(
            ln for ln in raw.splitlines() if not ln.strip().startswith("```")
        ).strip()
    return json.loads(raw)


def _discover_template(template_texts: list, doc_type: str, connection_id: str = None) -> dict:
    combined = "\n\n---\n\n".join(
        f"[Example {i+1}]\n{_trunc(t)}" for i, t in enumerate(template_texts)
    )
    return _llm_json([
        {
            "role": "system",
            "content": (
                "You are an expert in ISO 13485 software V&V documentation. "
                "Return only valid JSON."
            ),
        },
        {
            "role": "user",
            "content": (
                f'Analyze these example "{doc_type}" documents and extract their template structure.\n\n'
                f"{combined}\n\n"
                "Return JSON:\n"
                '{"sections":[{"heading":"...","description":"...","required_elements":["..."]}],'
                '"style_notes":"...","regulatory_language":["..."]}'
            ),
        },
    ], connection_id=connection_id)


def _research_inputs(input_texts: list, doc_type: str, connection_id: str = None) -> dict:
    combined = "\n\n---\n\n".join(
        f"[Input {i+1}]\n{_trunc(t)}" for i, t in enumerate(input_texts)
    )
    return _llm_json([
        {
            "role": "system",
            "content": (
                "You are an expert technical analyst for ISO 13485 regulated software. "
                "Return only valid JSON."
            ),
        },
        {
            "role": "user",
            "content": (
                f'Extract all information needed to write a "{doc_type}" from these input documents.\n\n'
                f"{combined}\n\n"
                "Return JSON:\n"
                '{"system_name":"...","system_description":"...","technical_specs":{},'
                '"requirements":[],"infrastructure":{},"stakeholders":[],'
                '"key_facts":[],"other_details":{}}'
            ),
        },
    ], connection_id=connection_id)


def _generate_section(
    doc_type: str, section: dict, research: dict,
    style: str, reg_lang: list, connection_id: str = None,
) -> str:
    required = "\n".join(f"- {e}" for e in section.get("required_elements", []))
    reg = "\n".join(f"- {p}" for p in reg_lang[:10])
    return _llm_call([
        {
            "role": "system",
            "content": (
                "You are an expert technical writer for ISO 13485 regulated "
                "software V&V documentation."
            ),
        },
        {
            "role": "user",
            "content": (
                f'Write the "{section["heading"]}" section for a "{doc_type}" '
                f"document per ISO 13485.\n\n"
                f"Purpose: {section['description']}\n"
                f"Required elements:\n{required}\n\n"
                f"System info:\n{json.dumps(research, indent=2)[:5000]}\n\n"
                f"Style: {style}\n"
                f"Regulatory language:\n{reg}\n\n"
                "Write professional prose only — do not repeat the section heading."
            ),
        },
    ], connection_id=connection_id)


def _assemble_docx(doc_type: str, sections: list) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    title = doc.add_heading(doc_type, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta = doc.add_paragraph(
        f"Standard: ISO 13485  |  Generated: {datetime.now().strftime('%Y-%m-%d')}"
    )
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph()

    for heading, content in sections:
        doc.add_heading(heading, level=1)
        for block in content.split("\n\n"):
            if block.strip():
                doc.add_paragraph(block.strip())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


print("✓ Cell 3 complete — LLM pipeline ready")


# ══════════════════════════════════════════════════════════════════════════════
# CELL 4 — UI  (run this cell to launch the app)
# ══════════════════════════════════════════════════════════════════════════════

# ── Application state ────────────────────────────────────────────────────────

_state = {
    "input_files": {},    # file_id -> {name, content_bytes}
    "template_files": {}, # file_id -> {name, content_bytes}
    "jobs": {},           # job_id -> {status, doc_type, progress, result_bytes, error}
}
_lock = threading.Lock()

# ── CSS injection ─────────────────────────────────────────────────────────────

display(HTML("""
<style>
.dg-app { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif; }

.dg-header {
    background: #1e293b;
    color: #f8fafc;
    border-radius: 10px 10px 0 0;
    padding: 16px 24px;
    display: flex;
    align-items: center;
    gap: 12px;
    margin-bottom: 0;
}
.dg-header h1 { font-size: 18px; font-weight: 700; margin: 0; color: #f8fafc; }
.dg-header p  { font-size: 12px; color: #94a3b8; margin: 0; }

.dg-body {
    background: #f1f5f9;
    border-radius: 0 0 10px 10px;
    padding: 20px;
    border: 1px solid #e2e8f0;
    border-top: none;
}

.dg-panel {
    background: #fff;
    border: 1px solid #e2e8f0;
    border-radius: 8px;
    padding: 20px;
    min-width: 0;
}

.dg-section-title {
    font-size: 14px;
    font-weight: 700;
    color: #1e293b;
    margin: 0 0 12px 0;
    padding-bottom: 8px;
    border-bottom: 2px solid #e2e8f0;
}

.dg-upload-label {
    font-size: 13px;
    font-weight: 600;
    color: #334155;
    margin-bottom: 4px;
}
.dg-hint { font-size: 12px; color: #64748b; margin-bottom: 8px; }

.dg-badge {
    display: inline-block;
    font-size: 11px;
    font-weight: 600;
    padding: 2px 7px;
    border-radius: 20px;
    margin-left: 6px;
    vertical-align: middle;
}
.dg-badge-blue   { background: #dbeafe; color: #1d4ed8; }
.dg-badge-purple { background: #ede9fe; color: #6d28d9; }

.dg-file-item {
    display: flex;
    align-items: center;
    gap: 8px;
    background: #f8fafc;
    border: 1px solid #e2e8f0;
    border-radius: 6px;
    padding: 5px 8px;
    margin-bottom: 4px;
    font-size: 13px;
}
.dg-file-name { flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; color: #334155; }
.dg-file-size { font-size: 11px; color: #94a3b8; flex-shrink: 0; }

.dg-divider { border: none; border-top: 1px solid #e2e8f0; margin: 16px 0; }

.dg-job {
    border: 1px solid #e2e8f0;
    border-radius: 8px;
    padding: 12px;
    margin-bottom: 10px;
    background: #f8fafc;
    font-size: 13px;
}
.dg-job-complete { border-color: #86efac; background: #f0fdf4; }
.dg-job-error    { border-color: #fca5a5; background: #fff5f5; }

.dg-job-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 6px; }
.dg-job-type   { font-weight: 600; color: #1e293b; }

.dg-status-badge {
    font-size: 11px;
    font-weight: 700;
    padding: 2px 8px;
    border-radius: 20px;
    text-transform: uppercase;
}
.dg-status-queued   { background: #e2e8f0; color: #475569; }
.dg-status-running  { background: #dbeafe; color: #1e40af; }
.dg-status-complete { background: #dcfce7; color: #15803d; }
.dg-status-error    { background: #fee2e2; color: #b91c1c; }

.dg-progress-text { font-size: 12px; color: #64748b; margin-top: 4px; }
.dg-error-text    { font-size: 12px; color: #dc2626; margin-top: 4px; }

.dg-empty {
    text-align: center;
    color: #94a3b8;
    font-size: 13px;
    padding: 24px 12px;
    border: 1px dashed #cbd5e1;
    border-radius: 8px;
}

.widget-button.dg-btn-generate {
    background: #2563eb !important;
    color: white !important;
    border: none !important;
    font-weight: 600 !important;
}
.widget-button.dg-btn-generate:disabled {
    background: #93c5fd !important;
    cursor: not-allowed !important;
}
.widget-button.dg-btn-download {
    background: #16a34a !important;
    color: white !important;
    border: none !important;
    font-weight: 600 !important;
    font-size: 12px !important;
}
</style>
"""))

# ── Widget helpers ────────────────────────────────────────────────────────────

def _human_size(n):
    if n < 1024:
        return f"{n} B"
    if n < 1048576:
        return f"{n/1024:.1f} KB"
    return f"{n/1048576:.1f} MB"

def _file_emoji(name):
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    return {"pdf": "📕", "docx": "📘", "doc": "📘",
            "xlsx": "📗", "xls": "📗", "pptx": "📙", "ppt": "📙"}.get(ext, "📄")

# ── Layout constants ──────────────────────────────────────────────────────────

PANEL_LAYOUT  = widgets.Layout(width="100%", padding="4px")
FULL_WIDTH    = widgets.Layout(width="100%")
UPLOAD_BTN    = widgets.Layout(width="100%", margin="4px 0")
OUTPUT_AREA   = widgets.Layout(width="100%", min_height="40px")

# ── Output widgets for dynamic areas ─────────────────────────────────────────

out_input_files    = widgets.Output(layout=OUTPUT_AREA)
out_template_files = widgets.Output(layout=OUTPUT_AREA)
out_jobs           = widgets.Output(layout=OUTPUT_AREA)

# ── File list renderer ────────────────────────────────────────────────────────

def _render_file_list(file_type: str):
    out = out_input_files if file_type == "input" else out_template_files
    with _lock:
        files = dict(_state[f"{file_type}_files"])

    out.clear_output(wait=True)
    with out:
        if not files:
            display(HTML('<div class="dg-hint" style="margin:4px 0">No files uploaded yet.</div>'))
            return
        rows = []
        for fid, info in files.items():
            icon = _file_emoji(info["name"])
            size = _human_size(len(info["content_bytes"]))
            remove_btn = widgets.Button(
                description="✕",
                layout=widgets.Layout(width="28px", height="26px", padding="0"),
                style=widgets.ButtonStyle(button_color="#fee2e2"),
            )

            def _make_handler(file_id, ft):
                def _handler(b):
                    with _lock:
                        _state[f"{ft}_files"].pop(file_id, None)
                    _render_file_list(ft)
                    _update_generate_btn()
                return _handler

            remove_btn.on_click(_make_handler(fid, file_type))
            row = widgets.HBox(
                [
                    widgets.HTML(
                        f'<div class="dg-file-item">'
                        f'{icon} <span class="dg-file-name" title="{info["name"]}">{info["name"]}</span>'
                        f'<span class="dg-file-size">{size}</span>'
                        f"</div>",
                        layout=widgets.Layout(flex="1"),
                    ),
                    remove_btn,
                ],
                layout=widgets.Layout(align_items="center", margin="0 0 4px 0"),
            )
            rows.append(row)
        for row in rows:
            display(row)

# ── Jobs renderer ─────────────────────────────────────────────────────────────

def _render_jobs():
    out_jobs.clear_output(wait=True)
    with _lock:
        jobs = {jid: dict(j) for jid, j in _state["jobs"].items()}

    with out_jobs:
        if not jobs:
            display(HTML('<div class="dg-empty">No documents generated yet.</div>'))
            return

        for jid, job in sorted(jobs.items(), key=lambda x: x[1].get("created_at", ""), reverse=True):
            status = job["status"]
            status_cls = f"dg-status-{status}"
            status_label = {
                "queued": "Queued", "running": "Running…",
                "complete": "Complete", "error": "Error",
            }.get(status, status)

            job_cls = "dg-job"
            if status == "complete":
                job_cls += " dg-job-complete"
            elif status == "error":
                job_cls += " dg-job-error"

            progress_html = ""
            if status in ("queued", "running"):
                progress_html = f'<div class="dg-progress-text">{job.get("progress", "")}</div>'
            elif status == "error":
                progress_html = f'<div class="dg-error-text">⚠ {job.get("error", "Unknown error")}</div>'

            display(HTML(
                f'<div class="{job_cls}">'
                f'  <div class="dg-job-header">'
                f'    <span class="dg-job-type">{job["doc_type"]}</span>'
                f'    <span class="dg-status-badge {status_cls}">{status_label}</span>'
                f"  </div>"
                f"  {progress_html}"
                f"</div>"
            ))

            if status == "complete":
                dl_btn = widgets.Button(
                    description="⬇ Download .docx",
                    layout=widgets.Layout(width="160px", height="30px"),
                )
                dl_btn.add_class("dg-btn-download")

                def _make_dl(job_id):
                    def _dl(b):
                        with _lock:
                            j = _state["jobs"].get(job_id)
                        if j and j.get("result_bytes"):
                            b64 = base64.b64encode(j["result_bytes"]).decode()
                            slug = j["doc_type"].lower().replace(" ", "_")
                            fname = f"{slug}_{datetime.now().strftime('%Y%m%d')}.docx"
                            display(Javascript(
                                f"var a=document.createElement('a');"
                                f"a.href='data:application/octet-stream;base64,{b64}';"
                                f"a.download='{fname}';"
                                f"document.body.appendChild(a);a.click();"
                                f"document.body.removeChild(a);"
                            ))
                    return _dl

                dl_btn.on_click(_make_dl(jid))
                display(dl_btn)

# ── Generate button state ─────────────────────────────────────────────────────

def _update_generate_btn():
    with _lock:
        has_input    = bool(_state["input_files"])
        has_template = bool(_state["template_files"])
    btn_generate.disabled = not (has_input and has_template and w_doc_type.value)

# ── File upload handler ────────────────────────────────────────────────────────

def _handle_upload(upload_widget, file_type: str, status_out: widgets.Output):
    files = _get_uploaded_files(upload_widget)
    if not files:
        return
    with status_out:
        clear_output(wait=True)
        display(HTML(f'<span style="color:#2563eb;font-size:12px">↑ Uploading {len(files)} file(s)…</span>'))

    added = 0
    for f in files:
        if f["content_bytes"]:
            fid = str(uuid.uuid4())
            with _lock:
                _state[f"{file_type}_files"][fid] = {
                    "name": f["name"],
                    "content_bytes": f["content_bytes"],
                }
            added += 1

    with status_out:
        clear_output(wait=True)
        if added:
            display(HTML(f'<span style="color:#16a34a;font-size:12px">✓ {added} file(s) added</span>'))

    _render_file_list(file_type)
    _update_generate_btn()
    upload_widget.value = {} if isinstance(upload_widget.value, dict) else ()

# ── Generation job runner ─────────────────────────────────────────────────────

def _run_job(job_id: str, doc_type: str, connection_id: str):
    def _set(status, **kw):
        with _lock:
            _state["jobs"][job_id].update({"status": status, **kw})
        _render_jobs()

    try:
        _set("running", progress="Extracting text from input documents…")
        with _lock:
            input_files    = list(_state["input_files"].values())
            template_files = list(_state["template_files"].values())

        input_texts    = [extract_text(f["name"], f["content_bytes"]) for f in input_files]
        template_texts = [extract_text(f["name"], f["content_bytes"]) for f in template_files]

        _set("running", progress="Analysing template documents…")
        structure = _discover_template(template_texts, doc_type, connection_id=connection_id)

        _set("running", progress="Deep-researching input documents…")
        research  = _research_inputs(input_texts, doc_type, connection_id=connection_id)

        style    = structure.get("style_notes", "")
        reg_lang = structure.get("regulatory_language", [])
        sections = structure.get("sections", [])
        total    = len(sections)
        sections_out = []

        for i, section in enumerate(sections, 1):
            _set("running", progress=f"Generating section {i}/{total}: {section['heading']}…")
            content = _generate_section(
                doc_type, section, research, style, reg_lang, connection_id=connection_id
            )
            sections_out.append((section["heading"], content))

        _set("running", progress="Assembling Word document…")
        docx_bytes = _assemble_docx(doc_type, sections_out)

        _set("complete", progress="Done", result_bytes=docx_bytes, error=None)

    except Exception as exc:
        logger.error("Job %s failed:\n%s", job_id, traceback.format_exc())
        _set("error", progress="Failed", error=str(exc))


# ── Widgets ───────────────────────────────────────────────────────────────────

# -- Input file upload --
w_input_upload = widgets.FileUpload(
    accept=".docx,.doc,.pdf,.xlsx,.xls,.pptx,.ppt",
    multiple=True,
    description="Add Files",
    layout=UPLOAD_BTN,
)
out_input_status = widgets.Output()

def _on_input_upload(change):
    _handle_upload(w_input_upload, "input", out_input_status)

w_input_upload.observe(_on_input_upload, names="value")

# -- Template file upload --
w_template_upload = widgets.FileUpload(
    accept=".docx,.doc,.pdf,.xlsx,.xls,.pptx,.ppt",
    multiple=True,
    description="Add Files",
    layout=UPLOAD_BTN,
)
out_template_status = widgets.Output()

def _on_template_upload(change):
    _handle_upload(w_template_upload, "template", out_template_status)

w_template_upload.observe(_on_template_upload, names="value")

# -- Doc type --
w_doc_type = widgets.Dropdown(
    options=[""] + DOC_TYPES,
    value="",
    description="",
    layout=FULL_WIDTH,
    style={"description_width": "0"},
)
w_doc_type.observe(lambda c: _update_generate_btn(), names="value")

# -- LLM connection --
w_llm_conn = widgets.Text(
    value="",
    placeholder=DEFAULT_LLM_CONNECTION_ID,
    description="",
    layout=FULL_WIDTH,
    style={"description_width": "0"},
)

# -- Generate button --
btn_generate = widgets.Button(
    description="⚡  Generate Document",
    disabled=True,
    layout=FULL_WIDTH,
)
btn_generate.add_class("dg-btn-generate")
out_generate_status = widgets.Output()

def _on_generate(b):
    out_generate_status.clear_output(wait=True)
    doc_type = w_doc_type.value
    if not doc_type:
        with out_generate_status:
            display(HTML('<span style="color:#dc2626;font-size:12px">⚠ Please select a document type.</span>'))
        return

    connection_id = w_llm_conn.value.strip() or None
    job_id = str(uuid.uuid4())

    with _lock:
        _state["jobs"][job_id] = {
            "status": "queued",
            "doc_type": doc_type,
            "progress": "Queued",
            "result_bytes": None,
            "error": None,
            "created_at": datetime.now().isoformat(),
        }

    _render_jobs()

    threading.Thread(
        target=_run_job,
        args=(job_id, doc_type, connection_id),
        daemon=True,
    ).start()

    with out_generate_status:
        display(HTML(
            f'<span style="color:#2563eb;font-size:12px">'
            f'✓ Job started for <b>{doc_type}</b>. '
            f'Progress shown in the Generated Documents panel below.'
            f'</span>'
        ))

btn_generate.on_click(_on_generate)

# ── Layout assembly ───────────────────────────────────────────────────────────

W_HALF = widgets.Layout(width="50%", padding="0 8px 0 0")
W_HALF_R = widgets.Layout(width="50%", padding="0 0 0 8px")

left_panel = widgets.VBox([
    widgets.HTML('<p class="dg-section-title">1. Upload Documents</p>'),

    widgets.HTML(
        '<p class="dg-upload-label">Input Documents'
        '<span class="dg-badge dg-badge-blue">Source files</span></p>'
        '<p class="dg-hint">Design specs, business docs, configs — .docx .pdf .xlsx .pptx</p>'
    ),
    w_input_upload,
    out_input_status,
    out_input_files,

    widgets.HTML('<hr class="dg-divider">'),

    widgets.HTML(
        '<p class="dg-upload-label">Template Examples'
        '<span class="dg-badge dg-badge-purple">Reference docs</span></p>'
        '<p class="dg-hint">Example documents showing the desired format and style</p>'
    ),
    w_template_upload,
    out_template_status,
    out_template_files,
], layout=widgets.Layout(width="50%", padding="0 10px 0 0"))


right_panel = widgets.VBox([
    widgets.HTML('<p class="dg-section-title">2. Generate Document</p>'),

    widgets.HTML('<p class="dg-upload-label">Document Type</p>'),
    w_doc_type,

    widgets.HTML(
        '<p class="dg-upload-label" style="margin-top:10px">LLM Connection ID '
        '<span style="font-weight:400;color:#64748b;font-size:12px">(leave blank for default)</span></p>'
    ),
    w_llm_conn,

    widgets.HTML('<div style="height:10px"></div>'),
    btn_generate,
    out_generate_status,

    widgets.HTML('<hr class="dg-divider">'),
    widgets.HTML('<p class="dg-section-title">3. Generated Documents</p>'),
    out_jobs,
], layout=widgets.Layout(width="50%", padding="0 0 0 10px"))


app = widgets.VBox([
    widgets.HTML("""
        <div class="dg-header dg-app">
            <span style="font-size:28px">📄</span>
            <div>
                <h1>V&amp;V Document Generator</h1>
                <p>ISO 13485 Software Verification &amp; Validation — Dataiku LLM Mesh</p>
            </div>
        </div>
    """),
    widgets.HBox(
        [left_panel, right_panel],
        layout=widgets.Layout(
            width="100%",
            padding="16px",
            background="#f1f5f9",
            border="1px solid #e2e8f0",
            border_radius="0 0 10px 10px",
        ),
    ),
], layout=widgets.Layout(width="100%", max_width="1200px"))

# ── Initial render & display ──────────────────────────────────────────────────

_render_file_list("input")
_render_file_list("template")
_render_jobs()
display(app)
